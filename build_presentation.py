#!/usr/bin/env python3
"""
Build a high-quality PipeRewind presentation (.pptx)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Constants ──────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)

# Colors
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)   # very dark navy
BG_CARD    = RGBColor(0x1A, 0x22, 0x3A)   # card backgrounds
ACCENT     = RGBColor(0x06, 0xB6, 0xD4)   # cyan/teal
ACCENT2    = RGBColor(0x22, 0xD3, 0xEE)   # lighter cyan
ACCENT_ORG = RGBColor(0xF9, 0x73, 0x16)   # orange accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY      = RGBColor(0xCB, 0xD5, 0xE1)   # light gray text
MGRAY      = RGBColor(0x94, 0xA3, 0xB8)   # medium gray
DGRAY      = RGBColor(0x47, 0x55, 0x69)   # dark gray

IMG_DIR = os.path.join(os.path.dirname(__file__), "images")
IMG_TECH = os.path.join(IMG_DIR, "technical")

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H


# ── Helpers ────────────────────────────────────────────────────────
def dark_bg(slide):
    """Fill slide background with dark navy."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_text(slide, left, top, width, height, text,
             font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=16, color=LGRAY, bold=False,
             indent=0, space_before=Pt(4), alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.level = indent
    p.space_before = space_before
    p.alignment = alignment
    return p

def accent_line(slide, top):
    """Draw a thin teal accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), top, Inches(12), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

def section_header(slide, title, subtitle=None):
    """Slide with large section title."""
    dark_bg(slide)
    accent_line(slide, Inches(3.1))
    add_text(slide, Inches(0.6), Inches(2.0), Inches(12), Inches(1.0),
             title, font_size=40, color=WHITE, bold=True,
             alignment=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(3.3), Inches(12), Inches(0.8),
                 subtitle, font_size=20, color=MGRAY, bold=False,
                 alignment=PP_ALIGN.LEFT)

def content_slide(slide, title):
    """Standard content slide with title and accent line."""
    dark_bg(slide)
    add_text(slide, Inches(0.6), Inches(0.3), Inches(11), Inches(0.7),
             title, font_size=30, color=ACCENT2, bold=True)
    accent_line(slide, Inches(1.05))

def card(slide, left, top, width, height, title, bullets, icon_text=None):
    """Draw a card with rounded rectangle background."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = RGBColor(0x33, 0x3D, 0x55)
    shape.line.width = Pt(1)

    tf = add_text(slide, left + Inches(0.2), top + Inches(0.15),
                  width - Inches(0.4), Inches(0.5),
                  title, font_size=17, color=ACCENT, bold=True)
    for b in bullets:
        add_para(tf, "• " + b, font_size=13, color=LGRAY)

    if icon_text:
        add_text(slide, left + Inches(0.2), top + height - Inches(0.45),
                 width - Inches(0.4), Inches(0.35),
                 icon_text, font_size=11, color=MGRAY)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 1: Title                                                ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
dark_bg(sl)

# Large title
add_text(sl, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.2),
         "PipeRewind", font_size=60, color=WHITE, bold=True,
         alignment=PP_ALIGN.LEFT)
# Subtitle
add_text(sl, Inches(0.8), Inches(2.9), Inches(11.5), Inches(0.8),
         "A Time-Travel Debugger for UNIX Shell Pipelines",
         font_size=26, color=ACCENT, bold=False, alignment=PP_ALIGN.LEFT)
# Accent line
accent_line(sl, Inches(3.8))
# Authors
add_text(sl, Inches(0.8), Inches(4.1), Inches(11.5), Inches(0.6),
         "Luiz Takahashi  &  Ehzoc Chavez",
         font_size=20, color=LGRAY, alignment=PP_ALIGN.LEFT)
add_text(sl, Inches(0.8), Inches(4.65), Inches(11.5), Inches(0.5),
         "Washington State University  •  Systems Programming  •  Spring 2026",
         font_size=15, color=MGRAY, alignment=PP_ALIGN.LEFT)
# Tag-line at bottom
add_text(sl, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5),
         "github.com/balloonjazz/Pipe-Rewind",
         font_size=13, color=DGRAY, alignment=PP_ALIGN.LEFT)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 2: Section — Elevator Pitch                              ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
section_header(sl, "Part 1: Elevator Pitch",
               "Project Goals  •  Course Themes  •  Architecture Overview")


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 3: The Problem                                           ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "The Problem: Pipelines Are Opaque")

tf = add_text(sl, Inches(0.6), Inches(1.3), Inches(6.5), Inches(5.5),
              "When a UNIX pipeline fails, you only see the final output.",
              font_size=18, color=WHITE, bold=True)
add_para(tf, "")
add_para(tf, "$ cat log.txt | grep ERROR | sort | uniq -c",
         font_size=16, color=ACCENT, bold=True)
add_para(tf, "(nothing)", font_size=15, color=ACCENT_ORG)
add_para(tf, "")
add_para(tf, "Where did the data disappear?", font_size=17, color=WHITE, bold=True)
add_para(tf, "• Was the file empty?", font_size=15, color=LGRAY)
add_para(tf, "• Did grep filter everything?", font_size=15, color=LGRAY)
add_para(tf, "• Did sort crash?", font_size=15, color=LGRAY)
add_para(tf, "")
add_para(tf, "Traditional workarounds:", font_size=17, color=WHITE, bold=True)
add_para(tf, "• tee — requires modifying the pipeline manually", font_size=15, color=LGRAY)
add_para(tf, "• strace — drowns you in thousands of syscalls", font_size=15, color=LGRAY)
add_para(tf, "• printf debugging — destroys pipeline transparency", font_size=15, color=LGRAY)

# Image on right
img_path = os.path.join(IMG_DIR, "pipeline_problem.png")
if os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(7.5), Inches(1.5), Inches(5.2), Inches(5.2))


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 4: The Solution                                         ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "The Solution: PipeRewind")

tf = add_text(sl, Inches(0.6), Inches(1.3), Inches(6.5), Inches(3.0),
              "PipeRewind sits between every pair of commands, recording every byte with precise timestamps.",
              font_size=18, color=WHITE)
add_para(tf, "")
add_para(tf, "Fully transparent — same output, zero code changes.",
         font_size=17, color=ACCENT, bold=True)
add_para(tf, "")
add_para(tf, "Record → Replay → Understand.", font_size=17, color=WHITE, bold=True)
add_para(tf, "• Captures DATA_IN, DATA_OUT, PROC_START, PROC_EXIT events",
         font_size=14, color=LGRAY)
add_para(tf, "• Nanosecond timestamps via CLOCK_MONOTONIC",
         font_size=14, color=LGRAY)
add_para(tf, "• Binary .prt trace format with random-access index",
         font_size=14, color=LGRAY)
add_para(tf, "• Interactive ncurses TUI replay with time-scrubbing",
         font_size=14, color=LGRAY)

img_path = os.path.join(IMG_DIR, "piperewind_solution.png")
if os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(7.5), Inches(1.5), Inches(5.2), Inches(5.2))


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 5: Course Themes                                        ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "Systems Programming Themes")

card(sl, Inches(0.5), Inches(1.4), Inches(3.8), Inches(2.5),
     "Process Control",
     ["fork() / execvp() to spawn stages",
      "waitpid() for lifecycle management",
      "SIGCHLD handling and exit codes"],
     "pipeline.c  •  capture.c")

card(sl, Inches(4.7), Inches(1.4), Inches(3.8), Inches(2.5),
     "Event-Driven Concurrency",
     ["epoll_create / epoll_ctl / epoll_wait",
      "Non-blocking O_NONBLOCK file descriptors",
      "Single-threaded multiplexing of N pipes"],
     "capture.c")

card(sl, Inches(8.9), Inches(1.4), Inches(3.8), Inches(2.5),
     "System-Level I/O",
     ["Raw read() / write() on kernel pipes",
      "dup2() for stdin/stdout redirection",
      "fwrite() / fseek() for binary trace I/O"],
     "capture.c  •  trace.c")

card(sl, Inches(0.5), Inches(4.2), Inches(3.8), Inches(2.5),
     "Process Scheduling Observation",
     ["/proc/[pid]/status polling",
      "Detect Running / Sleeping / Blocked",
      "Real-time state displayed in TUI"],
     "procstate.c  •  tui.c")

card(sl, Inches(4.7), Inches(4.2), Inches(3.8), Inches(2.5),
     "POSIX Threading",
     ["pthread_create for live-mode capture",
      "Shared trace file between threads",
      "fflush() synchronization strategy"],
     "tui.c (tui_run_live)")

card(sl, Inches(8.9), Inches(4.2), Inches(3.8), Inches(2.5),
     "Binary File Formats",
     ["Custom .prt header/event/index layout",
      "Random-access via terminal index table",
      "Dynamic index reallocation for live mode"],
     "trace.h  •  trace.c")


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 6: Architecture Overview                                 ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "System Architecture")

img_path = os.path.join(IMG_TECH, "process_hierarchy.png")
if os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), Inches(6.5), Inches(6.0))

tf = add_text(sl, Inches(7.3), Inches(1.4), Inches(5.5), Inches(5.5),
              "~2,900 lines of C99", font_size=18, color=WHITE, bold=True)
add_para(tf, "")
add_para(tf, "7 source files, each with a single responsibility:",
         font_size=15, color=LGRAY)
add_para(tf, "")
add_para(tf, "main.c — CLI dispatcher (record/replay/live/dump)",
         font_size=13, color=LGRAY)
add_para(tf, "pipeline.c — State-machine command parser",
         font_size=13, color=LGRAY)
add_para(tf, "capture.c — epoll capture engine",
         font_size=13, color=LGRAY)
add_para(tf, "trace.c — Binary trace reader/writer",
         font_size=13, color=LGRAY)
add_para(tf, "tui.c — ncurses interactive viewer",
         font_size=13, color=LGRAY)
add_para(tf, "diff.c — LCS-based side-by-side diff",
         font_size=13, color=LGRAY)
add_para(tf, "procstate.c — /proc state monitor",
         font_size=13, color=LGRAY)
add_para(tf, "")
add_para(tf, "Dependencies: GCC, Linux epoll, ncurses, pthreads",
         font_size=14, color=MGRAY)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 7: Section — Design Choices                              ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
section_header(sl, "Part 2: Design Choices & Trade-offs",
               "Key Decisions  •  Trade-offs  •  Limitations")


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 8: Design Decisions                                      ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "Key Design Decisions")

card(sl, Inches(0.5), Inches(1.4), Inches(5.9), Inches(2.6),
     "Why epoll instead of select/poll?",
     ["O(1) readiness notification vs O(N) scanning",
      "Scales to arbitrary pipeline depth",
      "Native Linux kernel support, zero dependencies",
      "Edge-triggered mode avoids redundant wakeups"],
     "capture.c:  epoll_ctl(EPOLLIN)")

card(sl, Inches(6.8), Inches(1.4), Inches(5.9), Inches(2.6),
     "Why a custom binary format instead of JSON/text?",
     ["Variable-length events need byte-level control",
      "Index table enables O(1) random access via fseek()",
      "10-100x smaller files than equivalent text logs",
      "Live append + reader refresh with no parsing overhead"],
     "trace.h:  TraceEvent / IndexEntry structs")

card(sl, Inches(0.5), Inches(4.3), Inches(5.9), Inches(2.6),
     "Why parent-process interposition?",
     ["No ptrace/LD_PRELOAD — no privilege escalation needed",
      "Transparent: child processes are unmodified",
      "Parent controls pipe topology with dup2()",
      "Clean signal handling and waitpid() from one place"],
     "capture.c:  fork/exec/dup2 orchestration")

card(sl, Inches(6.8), Inches(4.3), Inches(5.9), Inches(2.6),
     "Why a sliding window in the TUI?",
     ["100K+ event traces would exhaust RAM if fully loaded",
      "Window of 20K events keeps memory under ~16 MB",
      "Dynamic load/free on navigation boundaries",
      "Index table makes any window position O(1) to seek"],
     "tui.c:  tui_load_window()")


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 9: Trade-offs & Limitations                              ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "Trade-offs & Known Limitations")

tf = add_text(sl, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.8),
              "Trade-offs We Made", font_size=22, color=ACCENT, bold=True)
add_para(tf, "")
add_para(tf, "Single-threaded capture engine",
         font_size=16, color=WHITE, bold=True)
add_para(tf, "Pro: No locking, no race conditions, simpler code",
         font_size=13, color=LGRAY)
add_para(tf, "Con: CPU-bound stages could starve the epoll loop",
         font_size=13, color=ACCENT_ORG)
add_para(tf, "")
add_para(tf, "fflush() for live synchronization",
         font_size=16, color=WHITE, bold=True)
add_para(tf, "Pro: Simple, no mutexes needed between writer/reader",
         font_size=13, color=LGRAY)
add_para(tf, "Con: Slight I/O overhead from frequent disk flushes",
         font_size=13, color=ACCENT_ORG)
add_para(tf, "")
add_para(tf, "In-process shell parsing versus calling /bin/sh",
         font_size=16, color=WHITE, bold=True)
add_para(tf, "Pro: Full control, handles quotes/escapes precisely",
         font_size=13, color=LGRAY)
add_para(tf, "Con: No support for redirections, subshells, or variables",
         font_size=13, color=ACCENT_ORG)

tf2 = add_text(sl, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.8),
               "Known Limitations", font_size=22, color=ACCENT_ORG, bold=True)
add_para(tf2, "")
add_para(tf2, "• Linux-only (epoll, /proc are Linux-specific APIs)",
         font_size=14, color=LGRAY)
add_para(tf2, "")
add_para(tf2, "• No stderr capture — only stdout/stdin pipes are traced",
         font_size=14, color=LGRAY)
add_para(tf2, "")
add_para(tf2, "• Shell features like $VAR, &&, ||, >file not supported",
         font_size=14, color=LGRAY)
add_para(tf2, "")
add_para(tf2, "• No persistent bookmarks or annotations in trace files",
         font_size=14, color=LGRAY)
add_para(tf2, "")
add_para(tf2, "• TUI requires real terminal (no SSH without pty allocation)",
         font_size=14, color=LGRAY)
add_para(tf2, "")
add_para(tf2, "Future work:", font_size=16, color=WHITE, bold=True)
add_para(tf2, "• stderr interception", font_size=14, color=MGRAY)
add_para(tf2, "• Trace merging for distributed pipelines", font_size=14, color=MGRAY)
add_para(tf2, "• Compressed .prt format (zlib)", font_size=14, color=MGRAY)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 10: Section — Feature Walkthrough                        ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
section_header(sl, "Part 3: Feature Walkthrough",
               "Build  •  Record  •  Dump  •  Replay TUI  •  Diff  •  Live Mode")


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 11: Building & Running                                   ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "Feature 1: Building & Running")

tf = add_text(sl, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5),
              "Build (one command):", font_size=20, color=WHITE, bold=True)
add_para(tf, "")
add_para(tf, "$ make clean && make", font_size=20, color=ACCENT, bold=True)
add_para(tf, "")
add_para(tf, "This compiles 7 source files → one piperewind binary.",
         font_size=16, color=LGRAY)
add_para(tf, "Makefile links: -lncurses -lpthread",
         font_size=14, color=MGRAY)
add_para(tf, "")
add_para(tf, "Run the automated test suite:", font_size=20, color=WHITE, bold=True)
add_para(tf, "")
add_para(tf, "$ make test", font_size=20, color=ACCENT, bold=True)
add_para(tf, "")
add_para(tf, "Tests 5 different pipeline configurations:",
         font_size=16, color=LGRAY)
add_para(tf, "• 2-stage (echo | cat)", font_size=14, color=LGRAY)
add_para(tf, "• 3-stage (echo | sort | head)", font_size=14, color=LGRAY)
add_para(tf, "• 4-stage (echo | sort | uniq -c | sort -rn)", font_size=14, color=LGRAY)
add_para(tf, "• Single-quoted arguments", font_size=14, color=LGRAY)
add_para(tf, "• Double-quoted arguments", font_size=14, color=LGRAY)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 12: Record Command                                       ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "Feature 2: Recording a Pipeline")

tf = add_text(sl, Inches(0.6), Inches(1.3), Inches(12), Inches(2.5),
              "$ ./piperewind record -v \"echo -e 'banana\\napple\\ncherry' | sort | head -2\"",
              font_size=17, color=ACCENT, bold=True)
add_para(tf, "")
add_para(tf, "What happens under the hood:", font_size=19, color=WHITE, bold=True)

card(sl, Inches(0.5), Inches(2.9), Inches(3.7), Inches(4.0),
     "1. Parse",
     ["State-machine tokenizer splits",
      "the command string at pipes",
      "Handles 'single' and \"double\"",
      "quote escaping correctly"],
     "pipeline.c")

card(sl, Inches(4.6), Inches(2.9), Inches(3.7), Inches(4.0),
     "2. Launch",
     ["fork() each stage as child",
      "dup2() redirects stdin/stdout",
      "to interposed capture pipes",
      "execvp() replaces child image"],
     "capture.c")

card(sl, Inches(8.7), Inches(2.9), Inches(3.7), Inches(4.0),
     "3. Capture & Write",
     ["epoll_wait() multiplexes all fds",
      "read() data, timestamp it,",
      "write DATA_OUT + DATA_IN events",
      "Forward bytes to next stage stdin"],
     "capture.c → trace.c")


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 13: Dump                                                 ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "Feature 3: Inspecting a Trace (dump)")

tf = add_text(sl, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5),
              "$ ./piperewind dump trace.prt", font_size=20, color=ACCENT, bold=True)
add_para(tf, "")
add_para(tf, "=== PipeRewind Trace: trace.prt ===", font_size=15, color=WHITE, bold=True)
add_para(tf, "Stages: 3   Events: 14", font_size=15, color=LGRAY)
add_para(tf, "")
add_para(tf, "  Stage 0 [pid 9213]: echo -e 'banana\\napple\\ncherry'",
         font_size=14, color=LGRAY)
add_para(tf, "  Stage 1 [pid 9214]: sort", font_size=14, color=LGRAY)
add_para(tf, "  Stage 2 [pid 9215]: head -2", font_size=14, color=LGRAY)
add_para(tf, "")
add_para(tf, "--- Events ---", font_size=15, color=WHITE, bold=True)
add_para(tf, "[  0.07 ms] stage 0  PROC_START", font_size=13, color=LGRAY)
add_para(tf, "[  1.54 ms] stage 0  DATA_OUT  20 bytes  \"banana\\napple\\ncherry\\n\"",
         font_size=13, color=ACCENT)
add_para(tf, "[  1.61 ms] stage 1  DATA_IN   20 bytes  \"banana\\napple\\ncherry\\n\"",
         font_size=13, color=LGRAY)
add_para(tf, "[  2.44 ms] stage 1  DATA_OUT  20 bytes  \"apple\\nbanana\\ncherry\\n\"",
         font_size=13, color=ACCENT)
add_para(tf, "[  2.50 ms] stage 2  DATA_IN   20 bytes  \"apple\\nbanana\\ncherry\\n\"",
         font_size=13, color=LGRAY)
add_para(tf, "[  2.58 ms] stage 2  DATA_OUT  13 bytes  \"apple\\nbanana\\n\"",
         font_size=13, color=ACCENT)
add_para(tf, "[  2.80 ms] stage 2  PROC_EXIT exit_code=0",
         font_size=13, color=LGRAY)
add_para(tf, "")
add_para(tf, "Every byte accounted for. Every timestamp recorded.",
         font_size=17, color=WHITE, bold=True)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 14: Replay TUI                                           ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "Feature 4: Interactive Replay TUI")

img_path = os.path.join(IMG_DIR, "tui_mockup.png")
if os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), Inches(6.0), Inches(6.0))

tf = add_text(sl, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5),
              "$ ./piperewind replay trace.prt",
              font_size=18, color=ACCENT, bold=True)
add_para(tf, "")
add_para(tf, "Full-screen ncurses interface:", font_size=17, color=WHITE, bold=True)
add_para(tf, "")
add_para(tf, "Header — file name, stage count, event count",
         font_size=14, color=LGRAY)
add_para(tf, "Timeline — progress bar with elapsed time",
         font_size=14, color=LGRAY)
add_para(tf, "Stages — color-coded status per process",
         font_size=14, color=LGRAY)
add_para(tf, "Data panel — raw payload at current time",
         font_size=14, color=LGRAY)
add_para(tf, "")
add_para(tf, "Navigation:", font_size=17, color=WHITE, bold=True)
add_para(tf, "j/k — step forward/backward by one event", font_size=14, color=LGRAY)
add_para(tf, "J/K — jump by 10 events", font_size=14, color=LGRAY)
add_para(tf, "Home/End — jump to first/last event", font_size=14, color=LGRAY)
add_para(tf, "↑/↓ — select pipeline stage", font_size=14, color=LGRAY)
add_para(tf, "d — toggle diff view", font_size=14, color=LGRAY)
add_para(tf, "f — toggle follow mode", font_size=14, color=LGRAY)
add_para(tf, "h — toggle hex dump view", font_size=14, color=LGRAY)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 15: Diff                                                 ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "Feature 5: Side-by-Side Diff View")

tf = add_text(sl, Inches(0.6), Inches(1.3), Inches(5.5), Inches(5.5),
              "Press 'd' in the TUI", font_size=20, color=ACCENT, bold=True)
add_para(tf, "")
add_para(tf, "Compares DATA_OUT of Stage N with DATA_IN of Stage N+1.",
         font_size=16, color=WHITE)
add_para(tf, "")
add_para(tf, "Uses Longest Common Subsequence (LCS) algorithm:",
         font_size=16, color=WHITE, bold=True)
add_para(tf, "• Same logic as git diff", font_size=14, color=LGRAY)
add_para(tf, "• O(n·m) dynamic programming", font_size=14, color=LGRAY)
add_para(tf, "• Lines highlighted: green = added, red = removed", font_size=14, color=LGRAY)
add_para(tf, "")
add_para(tf, "Why this matters:", font_size=17, color=WHITE, bold=True)
add_para(tf, "• Instantly see how each stage transforms data", font_size=14, color=LGRAY)
add_para(tf, "• Pinpoint exactly where corruption occurs", font_size=14, color=LGRAY)
add_para(tf, "• No manual copy-paste between temp files", font_size=14, color=LGRAY)

# Example on right
tf2 = add_text(sl, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5),
               "Example: sort stage", font_size=18, color=WHITE, bold=True)
add_para(tf2, "")
add_para(tf2, "Stage 0 output:     Stage 1 output:", font_size=15, color=MGRAY)
add_para(tf2, "")
add_para(tf2, '  banana               apple', font_size=16, color=ACCENT_ORG)
add_para(tf2, '  apple                banana', font_size=16, color=ACCENT)
add_para(tf2, '  cherry               cherry', font_size=16, color=LGRAY)
add_para(tf2, "")
add_para(tf2, "Lines are reordered — sort did its job.",
         font_size=15, color=LGRAY)
add_para(tf2, "If sort had mangled data, you'd see it instantly.",
         font_size=15, color=WHITE)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 16: Epoll Deep Dive                                      ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "Feature 6: The epoll Capture Loop (Deep Dive)")

img_path = os.path.join(IMG_TECH, "epoll_architecture.png")
if os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(0.3), Inches(1.3), Inches(6.5), Inches(6.0))

tf = add_text(sl, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.8),
              "One loop handles all N stages:", font_size=18, color=WHITE, bold=True)
add_para(tf, "")
add_para(tf, "1. epoll_wait() blocks until any pipe has data",
         font_size=14, color=LGRAY)
add_para(tf, "2. Lookup: which stage does this fd belong to?",
         font_size=14, color=LGRAY)
add_para(tf, "3. read() into 64KB buffer (CAPTURE_BUF_SIZE)",
         font_size=14, color=LGRAY)
add_para(tf, "4. Record DATA_OUT event for stage i",
         font_size=14, color=LGRAY)
add_para(tf, "5. Record DATA_IN event for stage i+1",
         font_size=14, color=LGRAY)
add_para(tf, "6. write() to forward pipe → stage i+1 stdin",
         font_size=14, color=LGRAY)
add_para(tf, "7. If read() returns 0: PIPE_EOF, close fd",
         font_size=14, color=LGRAY)
add_para(tf, "")
add_para(tf, "Performance:", font_size=17, color=ACCENT, bold=True)
add_para(tf, "• O(1) per event — no scanning of all fds",
         font_size=14, color=LGRAY)
add_para(tf, "• 64KB buffer matches kernel pipe capacity",
         font_size=14, color=LGRAY)
add_para(tf, "• Zero-copy forwarding, minimal latency overhead",
         font_size=14, color=LGRAY)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 17: Process State Monitoring                              ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "Feature 7: Process State Monitoring")

tf = add_text(sl, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5),
              "How does PipeRewind know if a process is blocked?",
              font_size=20, color=WHITE, bold=True)
add_para(tf, "")
add_para(tf, "It reads the /proc/[pid]/status pseudo-file.",
         font_size=18, color=ACCENT)
add_para(tf, "")
add_para(tf, "$ cat /proc/12345/status | grep State",
         font_size=16, color=ACCENT, bold=True)
add_para(tf, "State:    S (sleeping)",
         font_size=16, color=LGRAY)
add_para(tf, "")
add_para(tf, "The kernel exposes real-time scheduler state through /proc:",
         font_size=16, color=WHITE)
add_para(tf, "• R = Running / Runnable", font_size=15, color=LGRAY)
add_para(tf, "• S = Sleeping (waiting for I/O or signal)", font_size=15, color=LGRAY)
add_para(tf, "• D = Disk sleep (uninterruptible, blocked on I/O)", font_size=15, color=LGRAY)
add_para(tf, "• Z = Zombie (exited but not yet reaped)", font_size=15, color=LGRAY)
add_para(tf, "")
add_para(tf, "TUI color coding: green = Running, yellow = Sleeping, red = Exited",
         font_size=16, color=WHITE, bold=True)
add_para(tf, "This immediately reveals bottleneck stages in the pipeline.",
         font_size=15, color=LGRAY)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 18: Sliding Window                                       ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "Feature 8: Sliding Window Cache")

img_path = os.path.join(IMG_TECH, "sliding_window.png")
if os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.8))

tf = add_text(sl, Inches(0.6), Inches(1.3), Inches(6.0), Inches(5.8),
              "Problem: Huge traces crash the TUI", font_size=20, color=ACCENT_ORG, bold=True)
add_para(tf, "")
add_para(tf, "A 100K-event trace with payloads could be 200+ MB",
         font_size=15, color=LGRAY)
add_para(tf, "Loading everything into RAM is not viable.",
         font_size=15, color=LGRAY)
add_para(tf, "")
add_para(tf, "Solution: Sliding window of 20,000 events", font_size=18, color=WHITE, bold=True)
add_para(tf, "")
add_para(tf, "• tui_load_window(ts, center_idx) loads a chunk",
         font_size=14, color=LGRAY)
add_para(tf, "• Payloads are calloc/free'd dynamically",
         font_size=14, color=LGRAY)
add_para(tf, "• When user navigates near edge: free old, load new",
         font_size=14, color=LGRAY)
add_para(tf, "• Index table gives O(1) fseek() to any event",
         font_size=14, color=LGRAY)
add_para(tf, "")
add_para(tf, "Result: ~16 MB max regardless of trace size.",
         font_size=16, color=ACCENT, bold=True)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 19: Live Mode                                            ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "Feature 9: Live-Attach Mode")

tf = add_text(sl, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.8),
              "$ ./piperewind live \"tail -f log | grep ERROR\"",
              font_size=17, color=ACCENT, bold=True)
add_para(tf, "")
add_para(tf, "Records and displays simultaneously.",
         font_size=18, color=WHITE, bold=True)
add_para(tf, "")
add_para(tf, "Architecture:", font_size=17, color=WHITE, bold=True)
add_para(tf, "• Background pthread runs capture_run()",
         font_size=14, color=LGRAY)
add_para(tf, "• Main thread runs the ncurses TUI",
         font_size=14, color=LGRAY)
add_para(tf, "• Shared: trace.prt file on disk",
         font_size=14, color=LGRAY)
add_para(tf, "")
add_para(tf, "Synchronization:", font_size=17, color=WHITE, bold=True)
add_para(tf, "• Writer calls fflush() after every event",
         font_size=14, color=LGRAY)
add_para(tf, "• Reader calls trace_reader_refresh() in the getch() loop",
         font_size=14, color=LGRAY)
add_para(tf, "• halfdelay(1) gives 100ms non-blocking input timeout",
         font_size=14, color=LGRAY)
add_para(tf, "")
add_para(tf, "Follow mode (press 'f'):", font_size=17, color=WHITE, bold=True)
add_para(tf, "Automatically scrolls to the latest event.",
         font_size=14, color=LGRAY)

tf2 = add_text(sl, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5),
               "Thread model:", font_size=18, color=WHITE, bold=True)
add_para(tf2, "")
add_para(tf2, "┌─────────────────┐", font_size=13, color=MGRAY)
add_para(tf2, "│   Main Thread   │", font_size=13, color=ACCENT)
add_para(tf2, "│  ncurses TUI    │", font_size=13, color=MGRAY)
add_para(tf2, "│  getch() loop   │", font_size=13, color=MGRAY)
add_para(tf2, "│  refresh reader │", font_size=13, color=MGRAY)
add_para(tf2, "└────────┬────────┘", font_size=13, color=MGRAY)
add_para(tf2, "         │ shared trace.prt", font_size=13, color=LGRAY)
add_para(tf2, "┌────────┴────────┐", font_size=13, color=MGRAY)
add_para(tf2, "│ Background pthread│", font_size=13, color=ACCENT_ORG)
add_para(tf2, "│  capture_run()  │", font_size=13, color=MGRAY)
add_para(tf2, "│  epoll → write  │", font_size=13, color=MGRAY)
add_para(tf2, "│  fflush() sync  │", font_size=13, color=MGRAY)
add_para(tf2, "└─────────────────┘", font_size=13, color=MGRAY)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 20: Binary Trace Format                                  ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(sl, "Feature 10: Binary Trace File Format (.prt)")

tf = add_text(sl, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5),
              "On-disk layout optimized for random access:",
              font_size=18, color=WHITE, bold=True)
add_para(tf, "")
add_para(tf, "┌──────────────────────────────────────────────────────┐",
         font_size=14, color=MGRAY)
add_para(tf, "│  FileHeader (magic 0x50525400, version, num_stages) │   Fixed 32 bytes",
         font_size=14, color=ACCENT)
add_para(tf, "├──────────────────────────────────────────────────────┤",
         font_size=14, color=MGRAY)
add_para(tf, "│  StageHeader × N  (command, pid, per-stage meta)    │   N × 280 bytes",
         font_size=14, color=LGRAY)
add_para(tf, "├──────────────────────────────────────────────────────┤",
         font_size=14, color=MGRAY)
add_para(tf, "│  Event₀  Event₁  Event₂  ...  EventK               │   Variable-length",
         font_size=14, color=ACCENT)
add_para(tf, "│  (each: 24-byte header + payload bytes)             │",
         font_size=14, color=LGRAY)
add_para(tf, "├──────────────────────────────────────────────────────┤",
         font_size=14, color=MGRAY)
add_para(tf, "│  IndexEntry₀  IndexEntry₁  ...  IndexEntryK        │   K × 16 bytes",
         font_size=14, color=ACCENT)
add_para(tf, "│  (timestamp_ns, file_offset → fseek() target)      │",
         font_size=14, color=LGRAY)
add_para(tf, "└──────────────────────────────────────────────────────┘",
         font_size=14, color=MGRAY)
add_para(tf, "")
add_para(tf, "Index table = O(1) random access. No sequential parsing needed.",
         font_size=17, color=WHITE, bold=True)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 21: Live Demo Prompt                                     ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
section_header(sl, "Live Demo",
               "Building  •  Recording  •  Dumping  •  Replaying  •  Live Mode")


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 22: Thank You                                            ║
# ╚══════════════════════════════════════════════════════════════════╝
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
accent_line(sl, Inches(4.0))

add_text(sl, Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2),
         "Thank You", font_size=52, color=WHITE, bold=True,
         alignment=PP_ALIGN.LEFT)
add_text(sl, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.6),
         "Luiz Takahashi  &  Ehzoc Chavez",
         font_size=22, color=LGRAY, alignment=PP_ALIGN.LEFT)
add_text(sl, Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.5),
         "github.com/balloonjazz/Pipe-Rewind",
         font_size=16, color=ACCENT, alignment=PP_ALIGN.LEFT)
add_text(sl, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
         "Questions?",
         font_size=28, color=ACCENT_ORG, bold=True, alignment=PP_ALIGN.LEFT)


# ── Save ──────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(__file__), "PipeRewind_Presentation.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
